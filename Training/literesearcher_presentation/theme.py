"""
Claude-style theme + reusable slide helpers for LiteResearcher presentation.
Palette inspired by claude.ai: warm cream background, deep ink text,
Claude orange (#D97757) as the single accent. Generous whitespace.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ---- Palette ----
CREAM      = RGBColor(0xF5, 0xF1, 0xED)   # background
CREAM_SOFT = RGBColor(0xFA, 0xF7, 0xF2)   # surface card
INK        = RGBColor(0x1F, 0x1F, 0x1F)   # primary text
INK_SOFT   = RGBColor(0x6B, 0x63, 0x56)   # secondary text
WARM_GRAY  = RGBColor(0x8A, 0x80, 0x70)
DIVIDER    = RGBColor(0xE8, 0xE0, 0xD5)
ACCENT     = RGBColor(0xD9, 0x77, 0x57)   # Claude orange
ACCENT_DK  = RGBColor(0xB8, 0x5F, 0x44)
INK_PANEL  = RGBColor(0xEB, 0xE4, 0xD7)   # subtle box bg
GREEN      = RGBColor(0x5B, 0x8C, 0x5A)   # for pass@4 / good
NAVY       = RGBColor(0x2E, 0x5C, 0x8A)

# ---- Fonts ----
F_HEAD = "Georgia"             # falls back to Times on no-serif systems
F_BODY = "Calibri"             # PowerPoint default sans
F_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def make_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs, bg=CREAM):
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s

# ----------------- text helpers -----------------
def add_text(slide, x, y, w, h, text, *, font=F_BODY, size=14, color=INK,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb

def add_paragraphs(slide, x, y, w, h, lines, *, font=F_BODY, size=14, color=INK,
                   bold=False, align=PP_ALIGN.LEFT, line_spacing=1.25,
                   space_after=6):
    """lines: list of (text, opts_dict_or_None). opts overrides defaults."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        p.space_after = Pt(opts.get("space_after", space_after))
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.color.rgb = opts.get("color", color)
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", False)
    return tb

# ----------------- shape helpers -----------------
def add_rect(slide, x, y, w, h, fill=CREAM_SOFT, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    if radius:
        # tighten corner radius
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
    sh.shadow.inherit = False
    return sh

def add_hline(slide, x, y, w, *, color=DIVIDER, weight=0.75):
    ln = slide.shapes.add_connector(1, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

def add_vline(slide, x, y, h, *, color=DIVIDER, weight=0.75):
    ln = slide.shapes.add_connector(1, x, y, x, y + h)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

# ----------------- page chrome -----------------
def page_header(slide, eyebrow, title, *, accent_bar=True):
    if accent_bar:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.6), Inches(0.55),
                                     Inches(0.32), Inches(0.07))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
    add_text(slide, Inches(0.6), Inches(0.70), Inches(11), Inches(0.32),
             eyebrow.upper(), font=F_BODY, size=10.5, color=ACCENT, bold=True)
    add_text(slide, Inches(0.6), Inches(1.05), Inches(11), Inches(0.7),
             title, font=F_HEAD, size=30, color=INK, bold=False)
    add_hline(slide, Inches(0.6), Inches(1.78), Inches(12.13))

def page_footer(slide, idx, total, eyebrow="LiteResearcher · Presentation"):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             eyebrow, font=F_BODY, size=9, color=INK_SOFT)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{idx}  /  {total}", font=F_BODY, size=9, color=INK_SOFT,
             align=PP_ALIGN.RIGHT)
